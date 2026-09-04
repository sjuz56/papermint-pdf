from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from pathlib import Path
from typing import List
import tempfile
import shutil
import subprocess
import uuid
import threading
import time
from zipfile import ZipFile, ZIP_DEFLATED

import fitz
from pypdf import PdfReader, PdfWriter
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptInches
from openpyxl import Workbook
from PIL import Image
import pdfplumber
import pytesseract
from weasyprint import HTML
from lxml import etree


BASE = Path(__file__).parent
TMP = BASE / "tmp"
TMP.mkdir(exist_ok=True)

app = FastAPI(title="PaperMint PDF Toolbox")

app.mount(
    "/static",
    StaticFiles(directory=BASE / "static"),
    name="static",
)


TOOLS = [
    ("merge", "Merge PDF", "Combine PDFs in the order you want."),
    ("split", "Split PDF", "Split a PDF into separate files or ranges."),
    ("compress", "Compress PDF", "Reduce PDF file size while preserving quality."),
    ("pdf-word", "PDF to Word", "Convert PDF to an editable Word document."),
    ("pdf-ppt", "PDF to PowerPoint", "Convert each PDF page to a PowerPoint slide."),
    ("pdf-excel", "PDF to Excel", "Extract detected tables into an XLSX workbook."),
    ("pdf-jpg", "PDF to JPG", "Render PDF pages as JPG images."),
    ("word-pdf", "Word to PDF", "Convert DOC/DOCX to PDF."),
    ("ppt-pdf", "PowerPoint to PDF", "Convert PPT/PPTX to PDF."),
    ("excel-pdf", "Excel to PDF", "Convert XLS/XLSX to PDF."),
    ("jpg-pdf", "JPG to PDF", "Combine images into a PDF."),
    ("sign", "Sign PDF", "Add a simple text signature to a PDF page."),
    ("watermark", "Watermark", "Add text watermark to every page."),
    ("rotate", "Rotate PDF", "Rotate every page by 90, 180 or 270 degrees."),
    ("html-pdf", "HTML to PDF", "Convert uploaded HTML into PDF."),
    ("unlock", "Unlock PDF", "Remove password protection when you know the password."),
    ("protect", "Protect PDF", "Encrypt a PDF with a password."),
    ("organize", "Organize PDF", "Reorder pages using a page list such as 3,1,2."),
    ("pdfa", "PDF to PDF/A", "Create an archival-style PDF copy."),
    ("repair", "Repair PDF", "Rewrite a damaged/readable PDF into a fresh file."),
    ("page-numbers", "Page numbers", "Add page numbers to every page."),
    ("scan-pdf", "Scan to PDF", "Convert phone scans/images into a PDF."),
    ("ocr", "OCR PDF", "Recognize text from scanned PDF pages."),
    ("compare", "Compare PDF", "Create a text difference report for two PDFs."),
    ("redact", "Redact PDF", "Search and permanently redact specified text."),
    ("crop", "Crop PDF", "Crop all pages by margins in millimeters."),
]


@app.get("/", response_class=HTMLResponse)
def home():
    return (BASE / "static" / "index.html").read_text(encoding="utf-8")


@app.get("/api/tools")
def tools():
    return [
        {"id": a, "name": b, "description": c}
        for a, b, c in TOOLS
    ]


def save_upload(upload: UploadFile) -> Path:
    suffix = Path(upload.filename or "").suffix

    path = TMP / f"{uuid.uuid4().hex}{suffix}"

    with path.open("wb") as f:
        shutil.copyfileobj(upload.file, f)

    return path


def outpath(ext: str) -> Path:
    return TMP / f"{uuid.uuid4().hex}{ext}"


def libreoffice_to_pdf(src: Path) -> Path:
    outdir = Path(tempfile.mkdtemp(dir=TMP))

    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(outdir),
        str(src),
    ]

    cp = subprocess.run(
        cmd,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=120,
    )

    candidates = list(outdir.glob("*.pdf"))

    if not candidates:
        raise HTTPException(
            500,
            f"LibreOffice conversion failed: {cp.stderr[-500:]}",
        )

    return candidates[0]


def remove_pdf2docx_blank_pages(docx_path: Path):
    cleaned = docx_path.with_name(
        docx_path.stem + "_clean.docx"
    )

    ns = {
        "w":
        "http://schemas.openxmlformats.org/"
        "wordprocessingml/2006/main"
    }

    with ZipFile(docx_path, "r") as zin:

        root = etree.fromstring(
            zin.read("word/document.xml")
        )

        paragraphs = []

        page_breaks = root.xpath(
            './/w:br[@w:type="page"]',
            namespaces=ns,
        )

        for br in page_breaks:

            node = br

            while node is not None:

                if node.tag == (
                    "{"
                    + ns["w"]
                    + "}p"
                ):
                    paragraphs.append(node)
                    break

                node = node.getparent()

        processed = set()

        for paragraph in paragraphs:

            marker = id(paragraph)

            if marker in processed:
                continue

            processed.add(marker)

            parent = paragraph.getparent()

            if parent is not None:
                parent.remove(paragraph)

        xml = etree.tostring(
            root,
            xml_declaration=True,
            encoding="UTF-8",
            standalone="yes",
        )

        with ZipFile(
            cleaned,
            "w",
            ZIP_DEFLATED,
        ) as zout:

            for item in zin.infolist():

                if item.filename == "word/document.xml":

                    zout.writestr(
                        item,
                        xml,
                    )

                else:

                    zout.writestr(
                        item,
                        zin.read(item.filename),
                    )

    cleaned.replace(docx_path)


def pdf_to_word_editable(src: Path) -> Path:

    print(
        "PDF2DOCX ENGINE STARTED",
        flush=True,
    )

    from pdf2docx import Converter

    out = outpath(".docx")

    converter = Converter(str(src))

    try:

        converter.convert(
            str(out)
        )

    finally:

        converter.close()

    print(
        "PDF2DOCX ENGINE FINISHED",
        flush=True,
    )

    remove_pdf2docx_blank_pages(out)

    print(
        "DOCX POSTPROCESS FINISHED",
        flush=True,
    )

    return out


PDF_WORD_JOBS = {}
PDF_WORD_LOCK = threading.Lock()
PDF_WORD_TTL = 60 * 60


def cleanup_pdf_word_jobs():

    now = time.time()

    with PDF_WORD_LOCK:

        expired = [
            job_id
            for job_id, job
            in PDF_WORD_JOBS.items()
            if now - job.get(
                "created",
                now,
            ) > PDF_WORD_TTL
        ]

        for job_id in expired:

            job = PDF_WORD_JOBS.pop(
                job_id,
                None,
            )

            if not job:
                continue

            for key in (
                "source",
                "output",
            ):

                value = job.get(key)

                if value:

                    try:

                        Path(value).unlink(
                            missing_ok=True
                        )

                    except Exception:
                        pass


def run_pdf_word_job(
    job_id: str,
    source: Path,
):

    try:

        with PDF_WORD_LOCK:

            PDF_WORD_JOBS[
                job_id
            ]["status"] = "processing"

        output = pdf_to_word_editable(
            source
        )

        with PDF_WORD_LOCK:

            PDF_WORD_JOBS[
                job_id
            ]["status"] = "done"

            PDF_WORD_JOBS[
                job_id
            ]["output"] = str(output)

    except Exception as e:

        print(
            f"PDF WORD JOB ERROR "
            f"{job_id}: {e}",
            flush=True,
        )

        with PDF_WORD_LOCK:

            if job_id in PDF_WORD_JOBS:

                PDF_WORD_JOBS[
                    job_id
                ]["status"] = "error"

                PDF_WORD_JOBS[
                    job_id
                ]["error"] = str(e)


@app.post("/api/pdf-word/start")
async def pdf_word_start(
    file: UploadFile = File(...)
):

    cleanup_pdf_word_jobs()

    if not file.filename:

        raise HTTPException(
            400,
            "No file uploaded.",
        )

    if (
        Path(file.filename)
        .suffix
        .lower()
        != ".pdf"
    ):

        raise HTTPException(
            400,
            "Please upload a PDF file.",
        )

    source = save_upload(file)

    job_id = uuid.uuid4().hex

    with PDF_WORD_LOCK:

        PDF_WORD_JOBS[job_id] = {
            "status": "queued",
            "created": time.time(),
            "source": str(source),
            "output": None,
            "error": None,
        }

    thread = threading.Thread(
        target=run_pdf_word_job,
        args=(
            job_id,
            source,
        ),
        daemon=True,
    )

    thread.start()

    return {
        "job_id": job_id,
        "status": "queued",
    }


@app.get("/api/pdf-word/status/{job_id}")
def pdf_word_status(job_id: str):

    cleanup_pdf_word_jobs()

    with PDF_WORD_LOCK:

        job = PDF_WORD_JOBS.get(
            job_id
        )

        if not job:

            raise HTTPException(
                404,
                "Conversion job not found.",
            )

        return {
            "job_id": job_id,
            "status": job["status"],
            "error": job.get("error"),
        }


@app.get("/api/pdf-word/download/{job_id}")
def pdf_word_download(job_id: str):

    with PDF_WORD_LOCK:

        job = PDF_WORD_JOBS.get(
            job_id
        )

        if not job:

            raise HTTPException(
                404,
                "Conversion job not found.",
            )

        if job["status"] == "error":

            raise HTTPException(
                500,
                job.get("error")
                or "Conversion failed.",
            )

        if job["status"] != "done":

            raise HTTPException(
                409,
                "Conversion is not finished yet.",
            )

        output = job.get("output")

    if (
        not output
        or not Path(output).exists()
    ):

        raise HTTPException(
            404,
            "Converted file no longer exists.",
        )

    return FileResponse(
        output,
        filename="converted.docx",
        media_type=(
            "application/vnd.openxmlformats-"
            "officedocument.wordprocessingml.document"
        ),
    )


@app.post("/api/convert")
async def convert(
    tool: str = Form(...),
    files: List[UploadFile] = File(default=[]),
    mode: str = Form("editable"),
    password: str = Form(""),
    text: str = Form(""),
    rotation: int = Form(90),
    pages: str = Form(""),
    margin: float = Form(10.0),
    page_number_start: int = Form(1),
    signature_page: int = Form(1),
    signature_x: float = Form(30),
    signature_y: float = Form(30),
):

    if tool != "html-pdf" and not files:

        raise HTTPException(
            400,
            "Upload at least one file.",
        )

    paths = [
        save_upload(f)
        for f in files
    ]

    try:

        if tool == "merge":

            writer = PdfWriter()

            for path in paths:

                reader = PdfReader(
                    str(path)
                )

                for page in reader.pages:

                    writer.add_page(page)

            out = outpath(".pdf")

            writer.write(str(out))

            return FileResponse(
                out,
                filename="merged.pdf",
            )


        if tool == "split":

            reader = PdfReader(
                str(paths[0])
            )

            indexes = []

            if pages.strip():

                for token in pages.split(","):

                    token = token.strip()

                    if "-" in token:

                        a, b = map(
                            int,
                            token.split(
                                "-",
                                1,
                            ),
                        )

                        indexes.extend(
                            range(
                                a - 1,
                                b,
                            )
                        )

                    elif token:

                        indexes.append(
                            int(token) - 1
                        )

            else:

                indexes = [0]

            writer = PdfWriter()

            for i in indexes:

                if 0 <= i < len(
                    reader.pages
                ):

                    writer.add_page(
                        reader.pages[i]
                    )

            out = outpath(".pdf")

            writer.write(str(out))

            return FileResponse(
                out,
                filename="split.pdf",
            )


        if tool == "compress":

            doc = fitz.open(
                paths[0]
            )

            out = outpath(".pdf")

            doc.save(
                out,
                garbage=4,
                deflate=True,
                clean=True,
            )

            return FileResponse(
                out,
                filename="compressed.pdf",
            )


        if tool == "pdf-word":

            out = pdf_to_word_editable(
                paths[0]
            )

            return FileResponse(
                out,
                filename="converted.docx",
            )


        if tool == "pdf-ppt":

            pdf = fitz.open(
                paths[0]
            )

            prs = Presentation()

            prs.slide_width = (
                PptInches(13.333)
            )

            prs.slide_height = (
                PptInches(7.5)
            )

            blank = prs.slide_layouts[6]

            for page in pdf:

                slide = prs.slides.add_slide(
                    blank
                )

                pix = page.get_pixmap(
                    matrix=fitz.Matrix(
                        2,
                        2,
                    ),
                    alpha=False,
                )

                img = outpath(".png")

                pix.save(img)

                slide.shapes.add_picture(
                    str(img),
                    0,
                    0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                )

            out = outpath(".pptx")

            prs.save(out)

            return FileResponse(
                out,
                filename="converted.pptx",
            )


        if tool == "pdf-excel":

            wb = Workbook()

            wb.remove(
                wb.active
            )

            with pdfplumber.open(
                paths[0]
            ) as pdf:

                for pi, page in enumerate(
                    pdf.pages,
                    1,
                ):

                    tables = (
                        page.extract_tables()
                    )

                    if not tables:

                        ws = wb.create_sheet(
                            f"Page {pi}"
                        )

                        text_content = (
                            page.extract_text()
                            or ""
                        )

                        for ri, line in enumerate(
                            text_content.splitlines(),
                            1,
                        ):

                            ws.cell(
                                ri,
                                1,
                                line,
                            )

                    else:

                        for ti, table in enumerate(
                            tables,
                            1,
                        ):

                            ws = wb.create_sheet(
                                f"P{pi} T{ti}"[:31]
                            )

                            for r, row in enumerate(
                                table,
                                1,
                            ):

                                for c, value in enumerate(
                                    row,
                                    1,
                                ):

                                    ws.cell(
                                        r,
                                        c,
                                        value,
                                    )

            out = outpath(".xlsx")

            wb.save(out)

            return FileResponse(
                out,
                filename="converted.xlsx",
            )


        if tool == "pdf-jpg":

            pdf = fitz.open(
                paths[0]
            )

            if len(pdf) == 1:

                pix = pdf[0].get_pixmap(
                    matrix=fitz.Matrix(
                        2,
                        2,
                    ),
                    alpha=False,
                )

                out = outpath(".jpg")

                pix.save(out)

                return FileResponse(
                    out,
                    filename="page-1.jpg",
                )

            import zipfile

            out = outpath(".zip")

            with zipfile.ZipFile(
                out,
                "w",
            ) as z:

                for i, page in enumerate(
                    pdf,
                    1,
                ):

                    img = outpath(".jpg")

                    page.get_pixmap(
                        matrix=fitz.Matrix(
                            2,
                            2,
                        ),
                        alpha=False,
                    ).save(img)

                    z.write(
                        img,
                        f"page-{i}.jpg",
                    )

            return FileResponse(
                out,
                filename="pages.zip",
            )


        if tool in {
            "word-pdf",
            "ppt-pdf",
            "excel-pdf",
        }:

            out = libreoffice_to_pdf(
                paths[0]
            )

            return FileResponse(
                out,
                filename="converted.pdf",
            )


        if tool in {
            "jpg-pdf",
            "scan-pdf",
        }:

            images = [
                Image.open(path)
                .convert("RGB")
                for path in paths
            ]

            out = outpath(".pdf")

            images[0].save(
                out,
                save_all=True,
                append_images=images[1:],
            )

            return FileResponse(
                out,
                filename="images.pdf",
            )


        if tool == "watermark":

            doc = fitz.open(
                paths[0]
            )

            watermark = (
                text
                or "WATERMARK"
            )

            for page in doc:

                rect = page.rect

                page.insert_text(
                    (
                        rect.width * 0.22,
                        rect.height * 0.52,
                    ),
                    watermark,
                    fontsize=34,
                    overlay=True,
                )

            out = outpath(".pdf")

            doc.save(out)

            return FileResponse(
                out,
                filename="watermarked.pdf",
            )


        if tool == "rotate":

            doc = fitz.open(
                paths[0]
            )

            for page in doc:

                page.set_rotation(
                    (
                        page.rotation
                        + rotation
                    )
                    % 360
                )

            out = outpath(".pdf")

            doc.save(out)

            return FileResponse(
                out,
                filename="rotated.pdf",
            )


        if tool == "protect":

            doc = fitz.open(
                paths[0]
            )

            out = outpath(".pdf")

            doc.save(
                out,
                encryption=fitz.PDF_ENCRYPT_AES_256,
                owner_pw=password or "owner",
                user_pw=password or "password",
                permissions=(
                    fitz.PDF_PERM_ACCESSIBILITY
                    | fitz.PDF_PERM_PRINT
                ),
            )

            return FileResponse(
                out,
                filename="protected.pdf",
            )


        if tool == "unlock":

            doc = fitz.open(
                paths[0]
            )

            if (
                doc.needs_pass
                and not doc.authenticate(
                    password
                )
            ):

                raise HTTPException(
                    400,
                    "Incorrect password.",
                )

            out = outpath(".pdf")

            doc.save(out)

            return FileResponse(
                out,
                filename="unlocked.pdf",
            )


        if tool == "organize":

            doc = fitz.open(
                paths[0]
            )

            order = [
                int(x.strip()) - 1
                for x in pages.split(",")
                if x.strip()
            ]

            if not order:

                raise HTTPException(
                    400,
                    "Enter page order, e.g. 3,1,2",
                )

            outdoc = fitz.open()

            for i in order:

                if 0 <= i < len(doc):

                    outdoc.insert_pdf(
                        doc,
                        from_page=i,
                        to_page=i,
                    )

            out = outpath(".pdf")

            outdoc.save(out)

            return FileResponse(
                out,
                filename="organized.pdf",
            )


        if tool in {
            "repair",
            "pdfa",
        }:

            doc = fitz.open(
                paths[0]
            )

            out = outpath(".pdf")

            doc.save(
                out,
                garbage=4,
                deflate=True,
                clean=True,
            )

            filename = (
                "repaired.pdf"
                if tool == "repair"
                else "archive-copy.pdf"
            )

            return FileResponse(
                out,
                filename=filename,
            )


        if tool == "page-numbers":

            doc = fitz.open(
                paths[0]
            )

            for i, page in enumerate(
                doc,
                page_number_start,
            ):

                rect = page.rect

                page.insert_text(
                    (
                        rect.width / 2 - 10,
                        rect.height - 18,
                    ),
                    str(i),
                    fontsize=10,
                )

            out = outpath(".pdf")

            doc.save(out)

            return FileResponse(
                out,
                filename="numbered.pdf",
            )


        if tool == "ocr":

            pdf = fitz.open(
                paths[0]
            )

            doc = Document()

            for i, page in enumerate(
                pdf
            ):

                pix = page.get_pixmap(
                    matrix=fitz.Matrix(
                        2,
                        2,
                    ),
                    alpha=False,
                )

                img = Image.frombytes(
                    "RGB",
                    [
                        pix.width,
                        pix.height,
                    ],
                    pix.samples,
                )

                recognized = (
                    pytesseract
                    .image_to_string(
                        img,
                        lang="eng",
                    )
                )

                doc.add_paragraph(
                    recognized
                )

                if i < len(pdf) - 1:
                    doc.add_page_break()

            out = outpath(".docx")

            doc.save(out)

            return FileResponse(
                out,
                filename="ocr.docx",
            )


        if tool == "compare":

            if len(paths) < 2:

                raise HTTPException(
                    400,
                    "Upload two PDFs to compare.",
                )

            import difflib

            text_a = "\n".join(
                page.get_text()
                for page
                in fitz.open(paths[0])
            )

            text_b = "\n".join(
                page.get_text()
                for page
                in fitz.open(paths[1])
            )

            diff = "\n".join(
                difflib.unified_diff(
                    text_a.splitlines(),
                    text_b.splitlines(),
                    fromfile="PDF A",
                    tofile="PDF B",
                    lineterm="",
                )
            )

            out = outpath(".txt")

            out.write_text(
                diff,
                encoding="utf-8",
            )

            return FileResponse(
                out,
                filename="comparison.txt",
            )


        if tool == "redact":

            if not text.strip():

                raise HTTPException(
                    400,
                    "Enter text to redact.",
                )

            doc = fitz.open(
                paths[0]
            )

            for page in doc:

                for rect in page.search_for(
                    text
                ):

                    page.add_redact_annot(
                        rect,
                        fill=(0, 0, 0),
                    )

                page.apply_redactions()

            out = outpath(".pdf")

            doc.save(out)

            return FileResponse(
                out,
                filename="redacted.pdf",
            )


        if tool == "crop":

            doc = fitz.open(
                paths[0]
            )

            m = margin * 72 / 25.4

            for page in doc:

                rect = page.rect

                page.set_cropbox(
                    fitz.Rect(
                        rect.x0 + m,
                        rect.y0 + m,
                        rect.x1 - m,
                        rect.y1 - m,
                    )
                )

            out = outpath(".pdf")

            doc.save(out)

            return FileResponse(
                out,
                filename="cropped.pdf",
            )


        if tool == "sign":

            doc = fitz.open(
                paths[0]
            )

            index = max(
                0,
                min(
                    signature_page - 1,
                    len(doc) - 1,
                ),
            )

            page = doc[index]

            page.insert_text(
                (
                    signature_x * 72 / 25.4,
                    signature_y * 72 / 25.4,
                ),
                text or "Signature",
                fontsize=18,
            )

            out = outpath(".pdf")

            doc.save(out)

            return FileResponse(
                out,
                filename="signed.pdf",
            )


        if tool == "html-pdf":

            if paths:

                html = paths[0].read_text(
                    encoding="utf-8",
                    errors="ignore",
                )

            else:

                html = text

            out = outpath(".pdf")

            HTML(
                string=html
            ).write_pdf(out)

            return FileResponse(
                out,
                filename="page.pdf",
            )


        raise HTTPException(
            400,
            "Unknown tool.",
        )

    except HTTPException:
        raise

    except Exception as e:

        raise HTTPException(
            500,
            f"Conversion failed: {e}",
        )
